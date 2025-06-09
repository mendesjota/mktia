"""Streamlit app para analisar fachadas de bares Brahma.

Este aplicativo permite enviar imagens, PDFs ou arquivos PPTX e realizar uma
analise preliminar usando as APIs da OpenAI ou Gemini. O objetivo e servir
como ponto de partida para treinar um modelo mais robusto.
"""

from io import BytesIO
from pathlib import Path
import base64

import streamlit as st
from pptx import Presentation
from PIL import Image, UnidentifiedImageError
import PyPDF2
import openai
import google.generativeai as genai
from google.api_core import exceptions # Importa exceções da API do Google

# --- Configurações da Página Streamlit ---
st.set_page_config(
    page_title="Oráculo de Fachadas Brahma",
    page_icon="🍺",
    layout="centered",
    initial_sidebar_state="expanded"
)

# --- Funções de Análise com APIs ---

def call_openai_api(prompt: str, api_key: str, image: Image.Image = None) -> str:
    """
    Envia o prompt (e opcionalmente uma imagem) para a API da OpenAI e retorna a resposta.
    Utiliza modelos de visão se uma imagem for fornecida.
    """
    if not api_key:
        return "Erro: Chave da API OpenAI não fornecida."

    openai.api_key = api_key
    messages_payload = []
    
    # Adiciona o texto do prompt como parte do conteúdo
    content_list = [{"type": "text", "text": prompt}]

    # Se uma imagem for fornecida, prepare-a para a API de visão
    if image:
        buffered = BytesIO()
        image.save(buffered, format="JPEG", quality=70) 
        img_str = base64.b64encode(buffered.getvalue()).decode("utf-8")
        content_list.append({"type": "image_url", "image_url": {"url": f"data:image:jpeg;base64,{img_str}", "detail": "low"}})
        model_name = "gpt-4o" # Modelo multimodal mais recente da OpenAI
    else:
        model_name = "gpt-3.5-turbo" # Modelo de texto padrão

    messages_payload.append({"role": "user", "content": content_list})

    try:
        response = openai.chat.completions.create(
            model=model_name,
            messages=messages_payload,
            max_tokens=500, # Limite de tokens para a resposta do LLM
        )
        return response.choices[0].message.content
    except openai.APIError as e:
        st.error(f"Erro da API OpenAI: {e.status_code} - {e.response.json().get('error', {}).get('message', 'Mensagem de erro não disponível')}")
        st.exception(e) # Mostra o erro completo para debug
        return f"O Oráculo encontrou um problema com a API OpenAI: {e.status_code}. Verifique sua chave e os limites de uso."
    except Exception as e:
        st.error(f"Ocorreu um erro inesperado com OpenAI: {e}")
        st.exception(e) # Mostra o erro completo para debug
        return "O Oráculo não conseguiu se comunicar com OpenAI. Tente novamente mais tarde."

def call_gemini_api(prompt: str, api_key: str, image: Image.Image = None) -> str:
    """
    Envia o prompt (e opcionalmente uma imagem) para a API Gemini e retorna a resposta.
    Utiliza modelos de visão se uma imagem for fornecida.
    """
    if not api_key:
        return "Erro: Chave da API Gemini não fornecida."

    try:
        genai.configure(api_key=api_key)
        
        if image:
            # Tenta usar o modelo de visão armazenado na session_state.
            # Se não existir ou for None, usa 'gemini-1.5-flash-latest' como fallback.
            model_name = st.session_state.get("gemini_vision_model", 'gemini-1.5-flash-latest')
            if model_name is None: # Verificação explícita se a validação falhou em encontrar um modelo
                return "Erro Gemini: Nenhum modelo de visão compatível encontrado. Por favor, valide sua chave API na barra lateral."
            
            model = genai.GenerativeModel(model_name) 
            response = model.generate_content([prompt, image])
        else:
            # Tenta usar o modelo de texto armazenado na session_state.
            # Se não existir ou for None, usa 'gemini-pro' como fallback.
            model_name = st.session_state.get("gemini_text_model", 'gemini-pro')
            if model_name is None: # Verificação explícita se a validação falhou em encontrar um modelo
                return "Erro Gemini: Nenhum modelo de texto compatível encontrado. Por favor, valide sua chave API na barra lateral."
            
            model = genai.GenerativeModel(model_name)
            response = model.generate_content(prompt)
            
        return response.text
    except exceptions.NotFound as e: # Captura o erro específico de modelo não encontrado
        st.error(f"Erro: O modelo Gemini '{model_name}' não foi encontrado. Por favor, valide sua chave API na barra lateral e verifique os modelos disponíveis para sua região.")
        st.exception(e) # Mostra o erro completo para debug
        return "O Oráculo não conseguiu acessar o modelo Gemini. Verifique sua chave API e os modelos disponíveis."
    except Exception as e:
        st.error(f"Ocorreu um erro inesperado com Gemini: {e}")
        st.exception(e) # Mostra o erro completo para debug
        return "O Oráculo não conseguiu se comunicar com Gemini. Tente novamente mais tarde. Verifique sua chave da API."

# --- Funções de Extração de Conteúdo ---

def analyze_content(prompt: str, provider: str, api_key: str, image: Image.Image = None) -> str:
    """
    Realiza a análise do conteúdo (texto ou imagem) usando o provedor escolhido.
    Esta função é um dispatcher para as APIs de OpenAI e Gemini.
    """
    if provider == "OpenAI":
        return call_openai_api(prompt, api_key, image=image)
    return call_gemini_api(prompt, api_key, image=image)

def extract_images_from_pptx(file_bytes: BytesIO) -> list[Image.Image]:
    """Extrai e retorna todas as imagens de um arquivo PPTX."""
    pres = Presentation(file_bytes)
    images: list[Image.Image] = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                try:
                    img_io = BytesIO(shape.image.blob)
                    img = Image.open(img_io)
                    # Converte para RGB para garantir compatibilidade com as APIs e evitar erros de modo
                    if img.mode in ("RGBA", "P"):
                        img = img.convert("RGB")
                    images.append(img)
                except UnidentifiedImageError:
                    st.warning(f"Um arquivo dentro do PPTX ({shape.name}) não pôde ser identificado como imagem. Pulando.")
                except Exception as e:
                    st.warning(f"Erro ao extrair imagem ({shape.name}) do PPTX: {e}. Pulando.")
    return images

def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extrai texto de um PDF. Para análise visual de PDF, considere PyMuPDF (fitz)."""
    try:
        reader = PyPDF2.PdfReader(BytesIO(file_bytes))
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text
    except PyPDF2.errors.PdfReadError:
        st.error("Erro ao ler o PDF. O arquivo pode estar corrompido ou protegido por senha.")
        return ""
    except Exception as e:
        st.error(f"Erro inesperado ao extrair texto do PDF: {e}")
        return ""

# --- Funções Auxiliares de UI ---
def display_example_images(sidebar_example_dir: Path):
    """Exibe imagens de exemplo na barra lateral, com tratamento de erros."""
    st.sidebar.markdown("---")
    st.sidebar.markdown("### Exemplos de Referência")
    
    if not sidebar_example_dir.exists():
        st.sidebar.warning(f"A pasta de exemplos '{sidebar_example_dir.name}' não foi encontrada.")
        st.sidebar.info(
            "Crie uma pasta chamada 'examples' na mesma pasta do seu script "
            "e adicione imagens JPEG ('.jpg' ou '.jpeg') lá."
        )
    else:
        # Pega todos os arquivos .jpg e .jpeg
        image_paths = sorted(list(sidebar_example_dir.glob("*.jpg")) + list(sidebar_example_dir.glob("*.jpeg")))
        
        if not image_paths:
            st.sidebar.info("Nenhuma imagem JPEG (.jpg ou .jpeg) encontrada na pasta 'examples'.")
        else:
            num_columns = min(3, len(image_paths))  # Mostrar no máximo 3 colunas
            columns = st.sidebar.columns(num_columns)
            
            for idx, image_path in enumerate(image_paths):
                try:
                    img_to_display = Image.open(image_path)
                    # Usa o módulo das colunas para exibir a imagem e a legenda
                    columns[idx % num_columns].image(img_to_display, caption=image_path.stem, use_column_width=True)
                except UnidentifiedImageError:
                    columns[idx % num_columns].error(
                        f"Não foi possível exibir '{image_path.name}'. "
                        "O arquivo pode estar corrompido ou não é uma imagem válida."
                    )
                except Exception as e:
                    columns[idx % num_columns].error(f"Erro inesperado ao carregar '{image_path.name}': {e}")


def validate_gemini_api_key(api_key: str):
    """Tenta configurar a API Gemini e lista os modelos disponíveis."""
    if not api_key:
        st.sidebar.warning("Por favor, digite sua chave da API Gemini.")
        return
    
    try:
        genai.configure(api_key=api_key)
        
        list_models_response = genai.list_models()
        
        found_text_model = None
        found_vision_model = None
        
        st.sidebar.success("Chave da API Gemini validada com sucesso! Modelos disponíveis:")
        
        model_names = []
        for m in list_models_response:
            model_names.append(m.name)
            if "generateContent" in m.supported_generation_methods:
                # Prioriza modelos mais recentes/multimodais para visão
                # Adicionado 'gemini-1.5-pro' aqui, pois é um modelo multimodal que pode ser usado com imagem
                if "gemini-1.5-flash" in m.name.lower() or "gemini-1.5-pro" in m.name.lower() or "vision" in m.name.lower():
                    if not found_vision_model:
                        found_vision_model = m.name
                # Prioriza gemini-pro para texto se não for um modelo de visão
                elif "gemini-pro" == m.name.lower(): 
                    if not found_text_model:
                        found_text_model = m.name
        
        st.sidebar.markdown(f"- **Texto:** `{found_text_model or 'Nenhum modelo de texto compatível encontrado.'}`")
        st.sidebar.markdown(f"- **Visão:** `{found_vision_model or 'Nenhum modelo de visão compatível encontrado.'}`")

        # Armazena na sessão para uso posterior, evitando re-listar modelos a cada re-run
        st.session_state["gemini_text_model"] = found_text_model
        st.session_state["gemini_vision_model"] = found_vision_model

        if not found_text_model and not found_vision_model:
            st.sidebar.error("Nenhum modelo Gemini compatível para 'generateContent' foi encontrado com esta chave de API. Verifique as permissões do seu projeto no Google Cloud.")

    except exceptions.FailedPrecondition as e:
        st.sidebar.error(f"Falha na validação da API Gemini: {e.message}. Verifique se a API do Gemini está ativada no seu projeto Google Cloud.")
        st.exception(e)
    except exceptions.InvalidArgument as e:
        st.sidebar.error(f"Falha na validação da API Gemini: {e.message}. A chave da API é inválida.")
        st.exception(e)
    except Exception as e:
        st.sidebar.error(f"Ocorreu um erro inesperado ao validar a API Gemini: {e}")
        st.exception(e)


# --- Interface Principal do Streamlit ---

def main() -> None:
    """Interface principal em Streamlit."""

    st.title("Oráculo de Fachadas Brahma 🍺")
    st.markdown(
        "<p style='font-size:1.1em; text-align:center;'> "
        "Envie imagens, PDFs ou apresentações PPTX para obter um diagnóstico inicial "
        "sobre fachadas Brahma. Informe sua chave da API na barra lateral para prosseguir."
        "</p>", unsafe_allow_html=True
    )
    st.markdown("---")

    # --- Sidebar (Configurações e Upload de Arquivos) ---
    with st.sidebar:
        st.header("Configurações do Oráculo")
        provider = st.selectbox("Escolha o Provedor da IA", ["OpenAI", "Gemini"], key="provider_select")
        api_key = st.text_input(f"Sua Chave de API ({provider})", type="password", key="api_key_input")

        if st.button("Validar API Key", key="validate_api_btn"):
            if provider == "OpenAI":
                if api_key:
                    openai.api_key = api_key
                    try:
                        # Tenta listar modelos para validar a chave OpenAI
                        openai.models.list() 
                        st.sidebar.success("Chave da API OpenAI validada com sucesso!")
                    except openai.APIError as e:
                        st.sidebar.error(f"Erro na API OpenAI: {e.status_code}. Verifique sua chave.")
                        st.exception(e)
                    except Exception as e:
                        st.sidebar.error(f"Erro inesperado ao validar OpenAI: {e}")
                        st.exception(e)
                else:
                    st.sidebar.warning("Por favor, digite sua chave da API OpenAI.")
            elif provider == "Gemini":
                validate_gemini_api_key(api_key)

        st.markdown("---")
        st.subheader("Envie seu Arquivo")
        uploaded_file = st.file_uploader(
            "Selecione uma imagem, PDF ou PPTX para análise:",
            type=["png", "jpg", "jpeg", "pdf", "pptx"],
            key="sidebar_file_uploader"
        )
        st.info(
            "Após enviar o arquivo, clique 'Enviar para Análise' no campo de chat "
            "para processá-lo e obter a resposta do Oráculo."
        )

        display_example_images(Path("examples"))

    # --- Histórico do Chat ---
    # Usando uma div HTML com altura fixa e overflow para o histórico de chat
    st.markdown("<div id='chat-history' style='height: 400px; overflow-y: auto; padding: 10px; border: 1px solid #ddd; border-radius: 8px;'>", unsafe_allow_html=True)
    
    if "messages" not in st.session_state:
        st.session_state.messages = []

    # Exibir Mensagens Anteriores
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            # Exibe imagem se existir
            if "image" in message and message["image"] is not None:
                if isinstance(message["image"], Image.Image):
                    st.image(message["image"], caption="Imagem enviada", width=200)
                elif isinstance(message["image"], (str, Path)): # Para exibir caminhos de arquivo (ex: imagens de exemplo)
                    st.image(str(message["image"]), caption="Imagem de Exemplo", width=200)
            
            # Exibe texto
            if "text" in message:
                st.markdown(message["text"])
            
            # Exibe preview de texto extraído se existir
            if "extracted_text_preview" in message and message["extracted_text_preview"]:
                with st.expander("Ver texto extraído do documento"):
                    st.text_area("Conteúdo extraído", message["extracted_text_preview"], height=150, disabled=True)
    
    st.markdown("</div>", unsafe_allow_html=True) # Fecha a div do histórico
    
    # --- Campo de Entrada do Chat (abaixo do histórico) ---
    # user_input_text é definido dentro do st.form e tem seu valor no submit
    # Aqui, garantimos que ele não será None se o formulário não foi submetido ainda
    user_input_text_form = "" # Variável auxiliar para o st.text_area

    with st.form("chat_form", clear_on_submit=True):
        user_input_text_form = st.text_area("O que você gostaria de perguntar ao Oráculo?", 
                                        height=80, 
                                        placeholder="Por exemplo: 'Analise esta fachada e me diga se ela segue o padrão Brahma.'",
                                        key="user_text_input") # Adicionado key para consistência
        
        submit_button = st.form_submit_button("Enviar para Análise")

    # A lógica de processamento deve ser acionada SOMENTE SE o botão foi clicado
    if submit_button:
        # Usamos user_input_text_form aqui, pois ele contém o valor do text_area no momento do submit
        user_input_text = user_input_text_form

        # uploaded_file é definido na sidebar, então está disponível aqui
        
        # Se um arquivo foi enviado na sidebar, mas nenhum texto foi digitado,
        # define um texto padrão para o prompt.
        if uploaded_file and not user_input_text:
            user_input_text = f"Analisar o arquivo: {uploaded_file.name}. Por favor, me dê um diagnóstico."

        if not api_key:
            st.error("Por favor, **informe a API Key** na barra lateral para iniciar a análise.")
            return # Sai da função main

        if not user_input_text and not uploaded_file:
            st.warning("Por favor, digite uma pergunta ou envie um arquivo para analisar.")
            return # Sai da função main

        # --- Lógica de Processamento da Mensagem do Usuário ---
        user_message_content = {"role": "user", "text": user_input_text}
        
        # Variável para a imagem que será enviada para a API (se houver)
        image_for_api_analysis = None 

        if uploaded_file:
            file_name = uploaded_file.name.lower()
            user_message_content["display_text"] = f"Você enviou o arquivo: `{file_name}`"
            
            if file_name.endswith((".png", ".jpg", ".jpeg")):
                try:
                    img = Image.open(uploaded_file)
                    user_message_content["image"] = img # Guarda a imagem PIL para exibir no chat
                    image_for_api_analysis = img # Define a imagem para ser analisada pela API
                except UnidentifiedImageError:
                    st.error("Erro: Não foi possível identificar a imagem enviada. Verifique se o arquivo está correto.")
                    return
                except Exception as e:
                    st.error(f"Erro ao carregar a imagem: {e}")
                    return
            elif file_name.endswith(".pptx"):
                user_message_content["file_type"] = "pptx"
                user_message_content["file_content"] = uploaded_file # Passa o BytesIO
            elif file_name.endswith(".pdf"):
                user_message_content["file_type"] = "pdf"
                user_message_content["file_content"] = uploaded_file.read() # Passa os bytes do PDF
            else:
                st.warning("Tipo de arquivo não suportado para análise direta. Por favor, use texto, imagem, PDF ou PPTX.")
                return

        # Adicionar a mensagem do usuário ao histórico antes de chamar a API
        st.session_state.messages.append(user_message_content)

        # Exibir a mensagem do usuário imediatamente
        with st.chat_message("user"):
            if "image" in user_message_content and user_message_content["image"] is not None:
                st.image(user_message_content["image"], caption="Imagem enviada", width=200)
            st.markdown(user_message_content.get("display_text", user_input_text))


        # --- Chamar a API e obter a resposta do Oráculo ---
        with st.chat_message("assistant"):
            with st.spinner("O Oráculo está consultando os ventos da análise..."):
                response_text = "Desculpe, o Oráculo está com problemas. Tente novamente mais tarde."
                extracted_text_for_chat = "" # Para armazenar texto extraído de PDF/PPTX para o histórico

                if image_for_api_analysis: # Se uma imagem foi carregada diretamente
                    # A IA analisa a imagem com o prompt inicial do usuário
                    response_text = analyze_content(user_input_text, provider, api_key, image=image_for_api_analysis)
                elif "file_type" in user_message_content and user_message_content["file_type"] == "pptx":
                    st.markdown("Extraindo imagens da apresentação...")
                    images_from_pptx = extract_images_from_pptx(user_message_content["file_content"])
                    if images_from_pptx:
                        full_analysis_output = []
                        for idx, img in enumerate(images_from_pptx):
                            st.image(img, caption=f"Imagem do Slide {idx + 1}", width=200)
                            # Analisa cada imagem do PPTX com um prompt específico para imagem
                            analysis_result = analyze_content(f"Analise a imagem {idx+1} da fachada:", provider, api_key, image=img)
                            full_analysis_output.append(f"**Análise da Imagem {idx + 1}:**\n{analysis_result}")
                        response_text = "\n\n".join(full_analysis_output)
                        if not full_analysis_output:
                             response_text = "Nenhuma imagem válida pôde ser extraída do PPTX para análise."
                    else:
                        response_text = "Não consegui extrair imagens desta apresentação PPTX para análise."
                elif "file_type" in user_message_content and user_message_content["file_type"] == "pdf":
                    st.markdown("Extraindo texto do PDF...")
                    text_from_pdf = extract_text_from_pdf(user_message_content["file_content"])
                    extracted_text_for_chat = text_from_pdf # Guarda para exibir no histórico
                    if text_from_pdf:
                        st.text_area("Texto extraído (para referência)", text_from_pdf, height=100, disabled=True)
                        # Analisa o texto extraído do PDF com o prompt do usuário
                        response_text = analyze_content(f"Analise o texto: {text_from_pdf}", provider, api_key)
                    else:
                        response_text = "Não consegui extrair texto deste PDF para análise."
                elif user_input_text: # Análise de texto puro
                    response_text = analyze_content(user_input_text, provider, api_key)
                
            st.markdown(response_text)
            # Adicionar a resposta do assistente ao histórico
            assistant_message_content = {"role": "assistant", "text": response_text}
            if extracted_text_for_chat:
                assistant_message_content["extracted_text_preview"] = extracted_text_for_chat
            st.session_state.messages.append(assistant_message_content)

# --- Execução do Aplicativo ---
if __name__ == "__main__":
    main()